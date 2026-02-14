#!/usr/bin/env python3
"""
Build the defense presentation using the university PPTX template.
Structure: 19 slides, ~15 minutes.
Uses pipeline_results.json for consistent numbers with figures.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import json
import os

REPO = "/home/user/Automotive-Sensors--Fault-Detection-and-Localization"
FIGS = os.path.join(REPO, "thesis_report", "figures")
TEMPLATE = os.path.join(REPO, "Temp.pptx")
OUTPUT = os.path.join(REPO, "thesis_report", "Defense_Presentation.pptx")

# Load pipeline results for consistent numbers
with open(os.path.join(REPO, "pipeline_results.json")) as f:
    R = json.load(f)

# Colors
DARK_BLUE = RGBColor(0, 51, 102)
MED_BLUE = RGBColor(0, 90, 160)
DARK_GREEN = RGBColor(0, 120, 60)
DARK_RED = RGBColor(180, 30, 30)
BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(255, 255, 255)
GRAY = RGBColor(80, 80, 80)
TABLE_HEADER_BG = RGBColor(0, 51, 102)
TABLE_ALT_BG = RGBColor(230, 240, 250)
HIGHLIGHT_BG = RGBColor(255, 255, 200)

prs = Presentation(TEMPLATE)

# Remove existing template slides
while len(prs.slides) > 0:
    sldId = prs.slides._sldIdLst[0]
    rId = sldId.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
    if rId is None:
        rId = sldId.get('r:id')
    if rId is not None:
        try:
            prs.part.drop_rel(rId)
        except KeyError:
            pass
    prs.slides._sldIdLst.remove(sldId)

title_layout = prs.slide_layouts[0]
object_layout = prs.slide_layouts[1]

# ---- Helper Functions ----
def add_slide(layout=None):
    return prs.slides.add_slide(layout or object_layout)

def add_textbox(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)

def set_text(tf, text, size=18, bold=False, color=BLACK, align=PP_ALIGN.LEFT, font_name='Calibri'):
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return p

def add_para(tf, text, size=16, bold=False, color=BLACK, align=PP_ALIGN.LEFT,
             space_before=Pt(4), space_after=Pt(2), font_name='Calibri', italic=False):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.font.italic = italic
    p.alignment = align
    p.space_before = space_before
    p.space_after = space_after
    return p

def add_bullet(tf, text, size=16, bold=False, color=BLACK, level=0, font_name='Calibri'):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.level = level
    p.space_before = Pt(3)
    p.space_after = Pt(2)
    return p

def add_image(slide, path, left, top, width=None, height=None):
    if os.path.exists(path):
        return slide.shapes.add_picture(path, left, top, width, height)
    return None

def set_cell(table, row, col, text, bold=False, color=BLACK, size=12):
    cell = table.cell(row, col)
    cell.text = str(text)
    for p in cell.text_frame.paragraphs:
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.size = Pt(size)
        p.font.name = 'Calibri'
        p.alignment = PP_ALIGN.CENTER

def style_table(table, header_bg=TABLE_HEADER_BG, font_size=12):
    for row_idx, row in enumerate(table.rows):
        for cell in row.cells:
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(font_size)
                p.font.name = 'Calibri'
                p.alignment = PP_ALIGN.CENTER
                if row_idx == 0:
                    p.font.bold = True
                    p.font.color.rgb = WHITE
            if row_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_bg
            elif row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_ALT_BG

# Layout constants
L = Inches(0.5)
T_TITLE = Inches(0.85)
T_CONTENT = Inches(1.55)
CONTENT_W = Inches(12.3)

def slide_title(slide, text):
    tb = add_textbox(slide, L, T_TITLE, CONTENT_W, Inches(0.55))
    tf = tb.text_frame; tf.word_wrap = True
    set_text(tf, text, size=26, bold=True, color=DARK_BLUE)

# Extract key numbers from pipeline results
t40 = R['thresholds']['40']
b40 = R['binary']['40']
pf40 = {pf['fault'].split('.')[0]: pf for pf in t40['per_fault']}

# Per-fault recall at 40th percentile
recall_acc_gain = pf40['acc fault gain']['recall']
recall_acc_noise = pf40['acc fault noise']['recall']
recall_acc_stuck = pf40['acc fault stuck']['recall']
recall_spd_gain = pf40['rpm fault gain']['recall']
recall_spd_noise = pf40['rpm fault noise']['recall']
recall_spd_stuck = pf40['rpm fault stuck at']['recall']

# =====================================================================
# SLIDE 1: TITLE
# =====================================================================
s = add_slide(title_layout)
for shape in list(s.shapes):
    if shape.has_text_frame:
        shape.text_frame.clear()

tb = add_textbox(s, Inches(0.5), Inches(0.6), Inches(9.5), Inches(0.5))
tf = tb.text_frame; tf.word_wrap = True
set_text(tf, "Institute for Software and Systems Engineering", size=14, color=GRAY)

tb = add_textbox(s, Inches(0.5), Inches(1.3), Inches(9.5), Inches(2.0))
tf = tb.text_frame; tf.word_wrap = True
set_text(tf, "Intelligent Analysis of Automotive Sensor Faults", size=28, bold=True, color=DARK_BLUE)
add_para(tf, "Using Self-Supervised Learning and Real-Time Simulation", size=24, bold=True, color=MED_BLUE, space_before=Pt(2))
add_para(tf, "", size=12)
add_para(tf, "Master Thesis Defense", size=18, color=GRAY, space_before=Pt(12))
add_para(tf, "", size=8)
add_para(tf, "Presented by:  Yahia Amir Yahia Gamal", size=16, color=BLACK, space_before=Pt(8))

tb = add_textbox(s, Inches(0.5), Inches(5.0), Inches(6.0), Inches(2.0))
tf = tb.text_frame; tf.word_wrap = True
set_text(tf, "First Examiner:     apl. Prof. Dr. Christoph Knieke", size=13, color=GRAY)
add_para(tf, "Second Examiner:  Dr. Stefan Wittek", size=13, color=GRAY)
add_para(tf, "Supervisor:            Dr. Mohammad Abboush", size=13, color=GRAY)
add_para(tf, "", size=8)
add_para(tf, "TU Clausthal  |  2025", size=14, bold=True, color=DARK_BLUE)

# =====================================================================
# SLIDE 2: OUTLINE
# =====================================================================
s = add_slide()
slide_title(s, "Outline")
tb = add_textbox(s, Inches(0.8), T_CONTENT, Inches(11.0), Inches(5.2))
tf = tb.text_frame; tf.word_wrap = True
items = [
    "Motivation and Problem Statement",
    "Research Questions",
    "Related Work: Why SimCLR?",
    "Proposed Framework (3 Phases)",
    "Datasets: A2D2 and HIL",
    "Fault Types and Injection Models",
    "SimCLR Architecture and 1D-CNN Encoder",
    "Anomaly Detection Method",
    "Evaluation Results",
    "Conclusion and Future Work",
]
set_text(tf, "", size=10)
for i, item in enumerate(items):
    add_bullet(tf, f"{i+1}.   {item}", size=18, color=BLACK)

# =====================================================================
# SLIDE 3: MOTIVATION
# =====================================================================
s = add_slide()
slide_title(s, "Motivation and Problem Statement")

tb = add_textbox(s, L, T_CONTENT, Inches(6.5), Inches(5.0))
tf = tb.text_frame; tf.word_wrap = True
set_text(tf, "The Problem:", size=20, bold=True, color=DARK_RED)
add_bullet(tf, "Modern vehicles contain 100+ sensors", size=15)
add_bullet(tf, "Sensor faults compromise safety-critical systems", size=15)
add_bullet(tf, "ISO 26262 requires rigorous validation", size=15)
add_bullet(tf, "Supervised methods need labeled fault data", size=15, bold=True)
add_bullet(tf, "Labeled data is expensive and inherently incomplete", size=15)
add_para(tf, "", size=8)
add_para(tf, "Our Solution:", size=20, bold=True, color=DARK_GREEN)
add_bullet(tf, "Learn only from healthy driving data", size=15, bold=True)
add_bullet(tf, "Self-supervised contrastive learning (SimCLR)", size=15)
add_bullet(tf, "Detect faults as deviations from normality", size=15)
add_bullet(tf, "Cross-domain: train on real data \u2192 test on HIL", size=15)

add_image(s, os.path.join(FIGS, "vmodel_diagram.png"),
          Inches(7.3), Inches(1.7), width=Inches(5.5))

# =====================================================================
# SLIDE 4: RESEARCH QUESTIONS
# =====================================================================
s = add_slide()
slide_title(s, "Research Questions")

shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.6), Inches(12.3), Inches(1.5))
shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor(230, 240, 255)
shape.line.color.rgb = MED_BLUE; shape.line.width = Pt(1.5)
tf = shape.text_frame; tf.word_wrap = True
set_text(tf, "RQ1: Can self-supervised contrastive learning on healthy data detect sensor faults?", size=16, bold=True, color=MED_BLUE)
add_para(tf, "Test core feasibility: train on A2D2 real driving data, detect faults in HIL simulation data.", size=14, color=BLACK)

shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(3.3), Inches(12.3), Inches(1.5))
shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor(230, 255, 240)
shape.line.color.rgb = DARK_GREEN; shape.line.width = Pt(1.5)
tf = shape.text_frame; tf.word_wrap = True
set_text(tf, "RQ2: How does the detection threshold affect precision-recall trade-off?", size=16, bold=True, color=DARK_GREEN)
add_para(tf, "Evaluate six threshold percentiles (15th\u201340th) to find optimal operating point.", size=14, color=BLACK)

shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5.0), Inches(12.3), Inches(1.5))
shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor(255, 240, 230)
shape.line.color.rgb = RGBColor(200, 100, 0); shape.line.width = Pt(1.5)
tf = shape.text_frame; tf.word_wrap = True
set_text(tf, "RQ3: Which fault types and sensors are most/least detectable?", size=16, bold=True, color=RGBColor(200, 100, 0))
add_para(tf, "Compare gain, noise, stuck-at faults on accelerator vs. speed sensors.", size=14, color=BLACK)

# =====================================================================
# SLIDE 5: RELATED WORK - WHY SimCLR? (NEW)
# =====================================================================
s = add_slide()
slide_title(s, "Related Work: Why SimCLR?")

add_image(s, os.path.join(FIGS, "ssl_comparison.png"),
          Inches(0.3), Inches(1.5), width=Inches(8.5), height=Inches(5.0))

tb = add_textbox(s, Inches(9.0), Inches(1.6), Inches(4.2), Inches(5.5))
tf = tb.text_frame; tf.word_wrap = True
set_text(tf, "Why SimCLR?", size=20, bold=True, color=DARK_BLUE)
add_para(tf, "", size=6)
add_bullet(tf, "Simplest architecture", size=15, bold=True, color=DARK_GREEN)
add_bullet(tf, "No momentum encoder (vs. MoCo)", size=14)
add_bullet(tf, "No asymmetric networks (vs. BYOL)", size=14)
add_bullet(tf, "Strong theory: NT-Xent loss", size=14)
add_para(tf, "", size=6)
add_para(tf, "Prior Applications:", size=17, bold=True, color=MED_BLUE)
add_bullet(tf, "Chen et al. (2020): ImageNet", size=13)
add_bullet(tf, "Tang et al. (2022): Bearing faults", size=13)
add_bullet(tf, "Wang et al. (2022): Motor diagnosis", size=13)
add_para(tf, "", size=6)
add_para(tf, "Gap: No automotive sensor", size=15, bold=True, color=DARK_RED)
add_para(tf, "fault detection with SimCLR", size=15, bold=True, color=DARK_RED)

# =====================================================================
# SLIDE 6: FRAMEWORK
# =====================================================================
s = add_slide()
slide_title(s, "Proposed Framework: Three-Phase Architecture")
add_image(s, os.path.join(FIGS, "framework_diagram.png"),
          Inches(0.3), Inches(1.5), width=Inches(12.5))

# Add small legend below
tb = add_textbox(s, Inches(0.5), Inches(5.8), Inches(12.0), Inches(1.5))
tf = tb.text_frame; tf.word_wrap = True
set_text(tf, "Phase 1 (green): Train SimCLR encoder on healthy A2D2 data  |  "
         "Phase 2 (blue): Calibrate with 90s healthy HIL data  |  "
         "Phase 3 (orange): Detect faults in HIL test data",
         size=13, color=GRAY)

# =====================================================================
# SLIDE 7: DATASETS
# =====================================================================
s = add_slide()
slide_title(s, "Datasets: A2D2 (Training) and HIL (Testing)")

tb = add_textbox(s, L, T_CONTENT, Inches(5.8), Inches(2.8))
tf = tb.text_frame; tf.word_wrap = True
set_text(tf, "A2D2 \u2014 Training Data (Phase 1)", size=17, bold=True, color=DARK_GREEN)
add_bullet(tf, f"Audi real-world driving, Munich", size=14)
add_bullet(tf, f"{R['a2d2']['total_samples']:,} samples ({R['a2d2']['duration_sec']/60:.1f} min)", size=14)
add_bullet(tf, "Sensors: Accelerator (%), Speed (km/h)", size=14)
add_bullet(tf, "Speed upsampled 50 Hz \u2192 100 Hz", size=14)
add_bullet(tf, "Only healthy data used for training", size=14, bold=True)

tb = add_textbox(s, Inches(6.8), T_CONTENT, Inches(5.8), Inches(2.8))
tf = tb.text_frame; tf.word_wrap = True
set_text(tf, "HIL \u2014 Test Data (Phases 2-3)", size=17, bold=True, color=DARK_RED)
add_bullet(tf, "dSPACE HIL simulator, TU Clausthal", size=14)
add_bullet(tf, "1 healthy + 6 fault recordings", size=14)
add_bullet(tf, "3 fault types \u00d7 2 sensors = 6 scenarios", size=14)
add_bullet(tf, f"Calibration: {R['healthy']['n_windows']} windows (90 sec)", size=14)
add_bullet(tf, "Cross-domain transfer: real \u2192 simulation", size=14, bold=True)

add_image(s, os.path.join(FIGS, "dataset_fusion_comparison.png"),
          Inches(1.5), Inches(4.2), width=Inches(10.0), height=Inches(3.0))

# =====================================================================
# SLIDE 8: FAULT TYPES
# =====================================================================
s = add_slide()
slide_title(s, "Fault Injection Types")
add_image(s, os.path.join(FIGS, "fault_types_comparison.png"),
          Inches(0.3), Inches(1.45), width=Inches(8.0), height=Inches(4.5))

tb = add_textbox(s, Inches(8.5), Inches(1.6), Inches(4.3), Inches(5.2))
tf = tb.text_frame; tf.word_wrap = True
set_text(tf, "Fault Models:", size=20, bold=True, color=DARK_BLUE)
add_para(tf, "", size=8)
add_para(tf, "1. Gain Fault:", size=17, bold=True, color=DARK_RED)
add_para(tf, "y(t) = \u03b1 \u00b7 x(t)", size=16, italic=True)
add_para(tf, "Calibration error, amplifier drift", size=13, color=GRAY)
add_para(tf, "", size=8)
add_para(tf, "2. Noise Fault:", size=17, bold=True, color=DARK_GREEN)
add_para(tf, "y(t) = x(t) + \u03b7(t),  \u03b7 ~ N(0,\u03c3\u00b2)", size=16, italic=True)
add_para(tf, "EMI, loose connections", size=13, color=GRAY)
add_para(tf, "", size=8)
add_para(tf, "3. Stuck-at Fault:", size=17, bold=True, color=RGBColor(125, 60, 152))
add_para(tf, "y(t) = c  (constant)", size=16, italic=True)
add_para(tf, "Hardware failure, frozen sensor", size=13, color=GRAY)

# =====================================================================
# SLIDE 9: SIMCLR + ENCODER
# =====================================================================
s = add_slide()
slide_title(s, "SimCLR Architecture and 1D-CNN Encoder")

add_image(s, os.path.join(FIGS, "simclr_architecture.png"),
          Inches(0.2), Inches(1.5), width=Inches(7.5), height=Inches(3.5))

# Key parameters on right
tb = add_textbox(s, Inches(7.8), Inches(1.5), Inches(5.2), Inches(3.5))
tf = tb.text_frame; tf.word_wrap = True
set_text(tf, "NT-Xent Loss (Contrastive):", size=17, bold=True, color=DARK_BLUE)
add_para(tf, "", size=4)
add_bullet(tf, "Maximize similarity: positive pair", size=14, color=DARK_GREEN)
add_bullet(tf, "Minimize similarity: 254 negatives", size=14, color=DARK_RED)
add_para(tf, "", size=6)
add_para(tf, "Training Configuration:", size=17, bold=True, color=MED_BLUE)
add_bullet(tf, "Temperature \u03c4 = 0.5", size=14)
add_bullet(tf, "Batch size = 128", size=14)
add_bullet(tf, "Epochs = 50", size=14)
add_bullet(tf, f"Windows: {R['training']['n_windows']:,}", size=14)
add_bullet(tf, f"Total parameters: {R['training']['total_params']:,}", size=14, bold=True)

# Encoder architecture below
add_image(s, os.path.join(FIGS, "encoder_architecture.png"),
          Inches(0.2), Inches(5.2), width=Inches(12.8), height=Inches(2.0))

# =====================================================================
# SLIDE 10: AUGMENTATION STRATEGIES
# =====================================================================
s = add_slide()
slide_title(s, "Data Preprocessing and Augmentation")

add_image(s, os.path.join(FIGS, "part2_augmentation_examples.png"),
          Inches(0.3), Inches(1.5), width=Inches(7.5), height=Inches(3.0))

tb = add_textbox(s, Inches(8.0), Inches(1.5), Inches(5.0), Inches(3.0))
tf = tb.text_frame; tf.word_wrap = True
set_text(tf, "3 Augmentation Strategies:", size=18, bold=True, color=DARK_BLUE)
add_para(tf, "", size=6)
add_bullet(tf, "Gaussian Jitter (\u03c3 = 0.1)", size=15)
add_para(tf, "  Simulates sensor noise", size=13, color=GRAY)
add_bullet(tf, "Amplitude Scaling [0.8, 1.2]", size=15)
add_para(tf, "  Simulates calibration variance", size=13, color=GRAY)
add_bullet(tf, "Temporal Masking (10%)", size=15)
add_para(tf, "  Forces temporal robustness", size=13, color=GRAY)

# Preprocessing info below
tb = add_textbox(s, Inches(0.5), Inches(4.8), Inches(12.3), Inches(2.5))
tf = tb.text_frame; tf.word_wrap = True
set_text(tf, "Preprocessing Pipeline:", size=18, bold=True, color=MED_BLUE)
add_para(tf, "", size=4)

# Create a nice summary table
table = s.shapes.add_table(4, 4, Inches(0.5), Inches(5.4), Inches(12.3), Inches(1.8)).table
headers = ["Step", "Method", "Parameters", "Output"]
for j, h in enumerate(headers):
    set_cell(table, 0, j, h, bold=True, color=WHITE, size=13)
data = [
    ["Normalization", "Z-Score", "\u03bc, \u03c3 from A2D2 training", "Zero mean, unit variance"],
    ["Segmentation", "Sliding Window", "W=200, stride=100 (50% overlap)", f"{R['training']['n_windows']:,} windows"],
    ["Augmentation", "Jitter + Scale + Mask", "Applied per window, per epoch", "2 views per window"],
]
for i, row in enumerate(data):
    for j, val in enumerate(row):
        set_cell(table, i+1, j, val, size=12)
style_table(table, font_size=12)

# =====================================================================
# SLIDE 11: ANOMALY DETECTION METHOD
# =====================================================================
s = add_slide()
slide_title(s, "Anomaly Detection Method")
add_image(s, os.path.join(FIGS, "anomaly_detection_flow.png"),
          Inches(0.3), Inches(1.5), width=Inches(12.5), height=Inches(2.5))

tb = add_textbox(s, Inches(0.5), Inches(4.2), Inches(5.8), Inches(3.0))
tf = tb.text_frame; tf.word_wrap = True
set_text(tf, "Phase 2: Calibration", size=18, bold=True, color=MED_BLUE)
add_bullet(tf, f"{R['healthy']['n_windows']} healthy HIL windows (90 s)", size=15)
add_bullet(tf, "Pass through frozen encoder", size=15)
add_bullet(tf, "Compute healthy centroid: \u03bc = mean(h)", size=15)
add_bullet(tf, "Set thresholds at 6 percentile levels", size=15)

tb = add_textbox(s, Inches(6.8), Inches(4.2), Inches(5.8), Inches(3.0))
tf = tb.text_frame; tf.word_wrap = True
set_text(tf, "Phase 3: Detection", size=18, bold=True, color=RGBColor(200, 100, 0))
add_bullet(tf, "Encode each test window \u2192 embedding h", size=15)
add_bullet(tf, "Compute: sim(h, \u03bc) = cosine similarity", size=15)
add_bullet(tf, "If sim < \u03c4 : FAULTY", size=15, bold=True, color=DARK_RED)
add_bullet(tf, "If sim \u2265 \u03c4 : HEALTHY", size=15, bold=True, color=DARK_GREEN)

# =====================================================================
# SLIDE 12: MULTI-THRESHOLD RESULTS TABLE
# =====================================================================
s = add_slide()
slide_title(s, "Results: Multi-Threshold Detection Performance")

add_image(s, os.path.join(FIGS, "part3_metrics_vs_threshold.png"),
          Inches(0.3), Inches(1.5), width=Inches(6.0), height=Inches(4.0))

table = s.shapes.add_table(7, 5, Inches(6.5), Inches(1.5), Inches(6.3), Inches(3.5)).table
headers = ["Percentile", "Threshold \u03c4", "Precision", "Recall", "F1-Score"]
for j, h in enumerate(headers):
    set_cell(table, 0, j, h, bold=True, color=WHITE, size=13)

for i, p in enumerate([15, 20, 25, 30, 35, 40]):
    td = R['thresholds'][str(p)]
    row_data = [
        f"{p}th",
        f"{td['threshold_value']:.3f}",
        f"{td['avg_precision']:.3f}",
        f"{td['avg_recall']:.3f}",
        f"{td['avg_f1']:.3f}",
    ]
    is_best = (p == 40)
    for j, val in enumerate(row_data):
        clr = DARK_GREEN if is_best else BLACK
        set_cell(table, i+1, j, val, bold=is_best, color=clr, size=13)

style_table(table, font_size=13)
# Highlight best row
for j in range(5):
    table.cell(6, j).fill.solid()
    table.cell(6, j).fill.fore_color.rgb = HIGHLIGHT_BG

tb = add_textbox(s, Inches(6.5), Inches(5.3), Inches(6.3), Inches(2.0))
tf = tb.text_frame; tf.word_wrap = True
set_text(tf, "Key Findings:", size=17, bold=True, color=DARK_BLUE)
add_bullet(tf, "Precision = 1.0 at ALL thresholds", size=15, bold=True, color=DARK_GREEN)
add_bullet(tf, f"Best F1 = {t40['avg_f1']:.3f} at 40th percentile", size=15, bold=True)
add_bullet(tf, f"ROC-AUC = {t40['roc_auc']:.3f}", size=15)

# =====================================================================
# SLIDE 13: PER-FAULT RECALL HEATMAP
# =====================================================================
s = add_slide()
slide_title(s, "Results: Per-Fault Recall Analysis")

add_image(s, os.path.join(FIGS, "part3_recall_heatmap.png"),
          Inches(0.3), Inches(1.5), width=Inches(6.5), height=Inches(4.0))

table = s.shapes.add_table(7, 4, Inches(7.0), Inches(1.5), Inches(5.8), Inches(3.2)).table
headers = ["Fault Type", "Detected", "Recall", "F1"]
for j, h in enumerate(headers):
    set_cell(table, 0, j, h, bold=True, color=WHITE, size=13)

fault_data = [
    ("Acc Gain", pf40['acc fault gain']),
    ("Acc Noise", pf40['acc fault noise']),
    ("Acc Stuck", pf40['acc fault stuck']),
    ("Speed Gain", pf40['rpm fault gain']),
    ("Speed Noise", pf40['rpm fault noise']),
    ("Speed Stuck", pf40['rpm fault stuck at']),
]
for i, (name, fd) in enumerate(fault_data):
    is_best = (fd['recall'] >= 0.99)
    is_worst = (fd['recall'] == min(f[1]['recall'] for f in fault_data))
    clr = DARK_GREEN if is_best else (DARK_RED if is_worst else BLACK)
    set_cell(table, i+1, 0, name, bold=is_best, color=clr, size=13)
    set_cell(table, i+1, 1, f"{fd['detected']}/{fd['windows']}", size=13)
    set_cell(table, i+1, 2, f"{fd['recall']:.1%}", bold=is_best, color=clr, size=13)
    set_cell(table, i+1, 3, f"{fd['f1']:.3f}", bold=is_best, color=clr, size=13)

style_table(table, font_size=13)
# Highlight best row (speed gain)
for j in range(4):
    table.cell(4, j).fill.solid()
    table.cell(4, j).fill.fore_color.rgb = HIGHLIGHT_BG

tb = add_textbox(s, Inches(7.0), Inches(5.0), Inches(5.8), Inches(2.3))
tf = tb.text_frame; tf.word_wrap = True
set_text(tf, "Observations (40th pctl):", size=17, bold=True, color=DARK_BLUE)
add_bullet(tf, f"Speed gain: {recall_spd_gain:.1%} recall (best)", size=14, bold=True, color=DARK_GREEN)
add_bullet(tf, f"Acc stuck: {recall_acc_stuck:.1%} recall (hardest)", size=14, color=DARK_RED)
add_bullet(tf, "Speed signals have lower natural variability", size=14)
add_bullet(tf, "Detectability reflects physical fault severity", size=14)

# =====================================================================
# SLIDE 14: BINARY CLASSIFICATION + CONFUSION MATRICES
# =====================================================================
s = add_slide()
slide_title(s, "Results: Binary Classification (Healthy vs. Faulty)")

add_image(s, os.path.join(FIGS, "part3_confusion_matrices.png"),
          Inches(0.1), Inches(1.45), width=Inches(7.5), height=Inches(4.5))

tb = add_textbox(s, Inches(7.8), Inches(1.5), Inches(5.2), Inches(2.0))
tf = tb.text_frame; tf.word_wrap = True
set_text(tf, "Binary Results (40th Percentile):", size=17, bold=True, color=DARK_BLUE)
add_para(tf, "", size=4)
add_bullet(tf, f"Accuracy:    {b40['accuracy']:.1%}", size=15, bold=True)
add_bullet(tf, f"Precision:    {b40['precision']:.1%}", size=15)
add_bullet(tf, f"Recall:        {b40['recall']:.1%}", size=15)
add_bullet(tf, f"F1-Score:    {b40['f1']:.3f}", size=15, bold=True)

# Confusion matrix table
table = s.shapes.add_table(3, 3, Inches(8.0), Inches(3.8), Inches(4.5), Inches(1.5)).table
set_cell(table, 0, 0, "", size=12)
set_cell(table, 0, 1, "Pred Healthy", bold=True, color=WHITE, size=12)
set_cell(table, 0, 2, "Pred Faulty", bold=True, color=WHITE, size=12)
set_cell(table, 1, 0, "True Healthy", bold=True, size=12)
set_cell(table, 1, 1, f"TN = {b40['tn']}", color=DARK_GREEN, size=14)
set_cell(table, 1, 2, f"FP = {b40['fp']}", color=DARK_RED, size=14)
set_cell(table, 2, 0, "True Faulty", bold=True, size=12)
set_cell(table, 2, 1, f"FN = {b40['fn']}", color=DARK_RED, size=14)
set_cell(table, 2, 2, f"TP = {b40['tp']:,}", bold=True, color=DARK_GREEN, size=14)
style_table(table)

tb = add_textbox(s, Inches(7.8), Inches(5.5), Inches(5.2), Inches(1.8))
tf = tb.text_frame; tf.word_wrap = True
set_text(tf, f"ROC-AUC = {t40['roc_auc']:.3f}", size=17, bold=True, color=MED_BLUE)
add_para(tf, "Significantly above random baseline (0.5)", size=14, color=GRAY)

# =====================================================================
# SLIDE 15: ROC + SIMILARITY DISTRIBUTIONS
# =====================================================================
s = add_slide()
slide_title(s, "ROC Curves and Similarity Distributions")

add_image(s, os.path.join(FIGS, "part3_roc_curves.png"),
          Inches(0.2), Inches(1.5), width=Inches(6.3), height=Inches(5.3))

add_image(s, os.path.join(FIGS, "part3_similarity_distributions.png"),
          Inches(6.5), Inches(1.5), width=Inches(6.3), height=Inches(4.0))

tb = add_textbox(s, Inches(6.5), Inches(5.7), Inches(6.3), Inches(1.5))
tf = tb.text_frame; tf.word_wrap = True
set_text(tf, "Clear separation between healthy and faulty distributions", size=15, bold=True, color=DARK_BLUE)
add_para(tf, "Higher threshold percentile \u2192 more recall, slightly more FP", size=14, color=GRAY)

# =====================================================================
# SLIDE 16: COMPUTING COSTS
# =====================================================================
s = add_slide()
slide_title(s, "Computing Cost Analysis")

table = s.shapes.add_table(7, 4, Inches(0.5), Inches(1.6), Inches(7.5), Inches(3.8)).table
headers = ["Phase", "Operation", "Time (s)", "% of Total"]
for j, h in enumerate(headers):
    set_cell(table, 0, j, h, bold=True, color=WHITE, size=14)

t = R['timing']
total = t['total']
data_rows = [
    ["Part 1", "A2D2 data loading", f"{t['part1']:.2f}", f"{100*t['part1']/total:.1f}%"],
    ["Part 2", "SimCLR training (50 epochs)", f"{t['training_only']:.2f}", f"{100*t['training_only']/total:.1f}%"],
    ["Part 2", "Augmentation + other", f"{t['part2_other']:.2f}", f"{100*t['part2_other']/total:.1f}%"],
    ["Part 3", "Embedding extraction", f"{t['part3_embeddings']:.2f}", f"{100*t['part3_embeddings']/total:.1f}%"],
    ["Part 3", "Evaluation + visualization", f"{t['part3_other']+t['part3_evaluation']:.2f}", f"{100*(t['part3_other']+t['part3_evaluation'])/total:.1f}%"],
    ["Total", "End-to-end pipeline", f"{total:.2f}", "100%"],
]
for i, row in enumerate(data_rows):
    for j, val in enumerate(row):
        bold = (i == 5)
        set_cell(table, i+1, j, val, bold=bold, size=14)
style_table(table, font_size=14)
# Highlight total row
for j in range(4):
    table.cell(6, j).fill.solid()
    table.cell(6, j).fill.fore_color.rgb = HIGHLIGHT_BG

tb = add_textbox(s, Inches(8.5), Inches(1.6), Inches(4.3), Inches(5.5))
tf = tb.text_frame; tf.word_wrap = True
set_text(tf, "Key Takeaways:", size=20, bold=True, color=DARK_BLUE)
add_para(tf, "", size=8)
add_bullet(tf, f"Training: {t['training_only']:.1f} s (one-time cost)", size=16, bold=True)
add_para(tf, "", size=6)
add_bullet(tf, "Inference: ~0.5 ms per window", size=16, bold=True)
add_para(tf, "", size=6)
add_bullet(tf, "Total pipeline: under 1 minute", size=16)
add_para(tf, "", size=6)
add_bullet(tf, "Suitable for real-time deployment", size=16)
add_para(tf, "", size=12)
add_para(tf, "vs. Supervised approach:", size=18, bold=True, color=DARK_RED)
add_para(tf, "  Ghannoum CNN-GRU: 23,000 s (6.4 h)", size=15, color=GRAY)
add_para(tf, "  Our approach: ~860x faster", size=16, bold=True, color=DARK_GREEN)
add_para(tf, "  + No labeled data required!", size=16, bold=True, color=DARK_GREEN)

# =====================================================================
# SLIDE 17: CONCLUSION
# =====================================================================
s = add_slide()
slide_title(s, "Conclusion")

# RQ1
shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(1.5), Inches(12.4), Inches(1.15))
shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor(230, 240, 255)
shape.line.color.rgb = MED_BLUE; shape.line.width = Pt(1)
tf = shape.text_frame; tf.word_wrap = True
set_text(tf, "RQ1:  Self-supervised representations from healthy A2D2 data successfully detect sensor faults in HIL data.", size=14, color=BLACK)
add_para(tf, "Cross-domain transfer works without any labeled fault examples. \u2713", size=14, bold=True, color=DARK_GREEN)

# RQ2
shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(2.85), Inches(12.4), Inches(1.15))
shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor(230, 255, 240)
shape.line.color.rgb = DARK_GREEN; shape.line.width = Pt(1)
tf = shape.text_frame; tf.word_wrap = True
set_text(tf, f"RQ2:  Per-fault precision = 1.0 at all thresholds. Recall: {R['thresholds']['15']['avg_recall']:.1%} \u2192 {t40['avg_recall']:.1%}.", size=14, color=BLACK)
add_para(tf, f"Optimal: 40th percentile (F1 = {t40['avg_f1']:.3f}). \u2713", size=14, bold=True, color=DARK_GREEN)

# RQ3
shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(4.2), Inches(12.4), Inches(1.15))
shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor(255, 240, 230)
shape.line.color.rgb = RGBColor(200, 100, 0); shape.line.width = Pt(1)
tf = shape.text_frame; tf.word_wrap = True
set_text(tf, f"RQ3:  Speed gain = {recall_spd_gain:.1%} recall (easiest). Acc stuck = {recall_acc_stuck:.1%} (hardest).", size=14, color=BLACK)
add_para(tf, "Detectability correlates with physical severity and sensor variability. \u2713", size=14, bold=True, color=RGBColor(200, 100, 0))

# Contributions
tb = add_textbox(s, Inches(0.5), Inches(5.55), Inches(12.3), Inches(1.8))
tf = tb.text_frame; tf.word_wrap = True
set_text(tf, "Key Contributions:", size=17, bold=True, color=DARK_BLUE)
add_bullet(tf, "First application of SimCLR to automotive sensor fault detection", size=14)
add_bullet(tf, "Successful cross-domain transfer: real driving data \u2192 HIL simulation", size=14)
add_bullet(tf, "Comprehensive multi-threshold evaluation with 6 fault scenarios", size=14)
add_bullet(tf, "Complete pipeline in under 1 minute on CPU (no GPU required)", size=14)

# =====================================================================
# SLIDE 18: FUTURE WORK
# =====================================================================
s = add_slide()
slide_title(s, "Future Work")

tb = add_textbox(s, Inches(0.5), T_CONTENT, Inches(5.8), Inches(5.0))
tf = tb.text_frame; tf.word_wrap = True
set_text(tf, "Short-term Extensions:", size=20, bold=True, color=DARK_BLUE)
add_para(tf, "", size=6)
add_bullet(tf, "Add more sensor channels", size=16)
add_para(tf, "  (RPM, steering, brake pressure)", size=13, color=GRAY)
add_bullet(tf, "Fault localization", size=16)
add_para(tf, "  Identify which sensor is faulty", size=13, color=GRAY)
add_bullet(tf, "Adaptive thresholding", size=16)
add_para(tf, "  Eliminate manual threshold selection", size=13, color=GRAY)
add_bullet(tf, "Test with real vehicle fault data", size=16)

tb = add_textbox(s, Inches(6.8), T_CONTENT, Inches(5.8), Inches(5.0))
tf = tb.text_frame; tf.word_wrap = True
set_text(tf, "Long-term Research:", size=20, bold=True, color=DARK_BLUE)
add_para(tf, "", size=6)
add_bullet(tf, "Compare SSL methods", size=16)
add_para(tf, "  TS2Vec, MoCo, BYOL benchmarking", size=13, color=GRAY)
add_bullet(tf, "Explainability (XAI)", size=16)
add_para(tf, "  Diagnostic insights for faults", size=13, color=GRAY)
add_bullet(tf, "Edge deployment", size=16)
add_para(tf, "  Real-time in-vehicle detection", size=13, color=GRAY)
add_bullet(tf, "ISO 26262 ASIL integration", size=16)
add_para(tf, "  Safety integrity level mapping", size=13, color=GRAY)

# =====================================================================
# SLIDE 19: THANK YOU
# =====================================================================
s = add_slide(title_layout)
for shape in list(s.shapes):
    if shape.has_text_frame:
        shape.text_frame.clear()

tb = add_textbox(s, Inches(1.0), Inches(1.5), Inches(11.0), Inches(5.0))
tf = tb.text_frame; tf.word_wrap = True
set_text(tf, "Thank You", size=44, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)
add_para(tf, "", size=16)
add_para(tf, "Questions?", size=32, color=MED_BLUE, align=PP_ALIGN.CENTER)
add_para(tf, "", size=20)
add_para(tf, "Yahia Amir Yahia Gamal", size=20, color=BLACK, align=PP_ALIGN.CENTER)
add_para(tf, "Institute of Software and Systems Engineering", size=16, color=GRAY, align=PP_ALIGN.CENTER)
add_para(tf, "TU Clausthal  |  2025", size=16, color=GRAY, align=PP_ALIGN.CENTER)

# ---- SAVE ----
prs.save(OUTPUT)
print(f"\nPresentation saved: {OUTPUT}")
print(f"Total slides: {len(prs.slides)}")
print("\nSlide structure:")
slide_titles = [
    "Title", "Outline", "Motivation", "Research Questions",
    "Related Work: Why SimCLR?", "Proposed Framework",
    "Datasets", "Fault Types", "SimCLR + Encoder",
    "Augmentation & Preprocessing", "Anomaly Detection",
    "Multi-Threshold Results", "Per-Fault Results",
    "Binary Classification", "ROC + Similarity",
    "Computing Costs", "Conclusion", "Future Work", "Thank You"
]
for i, t in enumerate(slide_titles):
    print(f"  Slide {i+1:2d}: {t}")
